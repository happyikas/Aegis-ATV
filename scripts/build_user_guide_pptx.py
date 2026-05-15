"""Generate ``docs/exports/USER_GUIDE.ko.pptx`` from the content of
``docs/USER_GUIDE.ko.md``.

This is a hand-curated distillation of the 13-section Korean user
guide into ~14 slides. Tone + palette inherits the NVIDIA Inception
PitchDeck: navy dominant, ice-blue secondary, coral accent for
warnings / big stats, all-Korean primary copy.

Run with::

    uv run --with python-pptx python scripts/build_user_guide_pptx.py

Output is committed as a binary asset (docs/exports/USER_GUIDE.ko.pptx).
Regenerate by re-running this script after content changes — the deck
is deterministic given the same script input.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── palette (PitchDeck-aligned: Midnight Executive + coral accent) ──
NAVY = RGBColor(0x1E, 0x27, 0x61)         # primary
ICE = RGBColor(0xCA, 0xDC, 0xFC)          # secondary
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CORAL = RGBColor(0xF9, 0x61, 0x67)        # warnings / "BLOCK"
MUTED = RGBColor(0x8A, 0x90, 0xA8)        # body subtle
SOFT_WHITE = RGBColor(0xF6, 0xF7, 0xFB)   # card backgrounds
DARK_INK = RGBColor(0x1A, 0x1E, 0x36)     # body text on light

# Korean-friendly font stack — Pretendard / Apple SD Gothic Neo / Malgun
HEAD_FONT = "Pretendard"
BODY_FONT = "Pretendard"
MONO_FONT = "Consolas"

# Slide size — 16:9 widescreen (default 13.33 × 7.5 in)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Default positions used as function defaults — pre-computed at module
# load to satisfy ruff B008 (no function calls in argument defaults).
DEFAULT_EYEBROW_TOP = Inches(0.42)
DEFAULT_TITLE_TOP = Inches(0.75)
DEFAULT_DIVIDER_TOP = Inches(1.55)


# ── helpers ─────────────────────────────────────────────────────


def _fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _outline(shape, color: RGBColor, width_pt: float = 0.75) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def _add_text(
    slide,
    text: str,
    *,
    left: Inches,
    top: Inches,
    width: Inches,
    height: Inches,
    size: int = 16,
    bold: bool = False,
    color: RGBColor = DARK_INK,
    font: str = BODY_FONT,
    align: str = "left",
    italic: bool = False,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    tb.text_frame.margin_left = Emu(0)
    tb.text_frame.margin_right = Emu(0)
    tb.text_frame.margin_top = Emu(0)
    tb.text_frame.margin_bottom = Emu(0)
    p = tb.text_frame.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def _add_bg(slide, color: RGBColor) -> None:
    """Solid background rectangle covering the whole slide."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H,
    )
    _fill(bg, color)


def _footer(slide, page_num: str, total: str = "15") -> None:
    """Slim navy footer band with page number."""
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.18), SLIDE_W, Inches(0.32),
    )
    _fill(band, NAVY)
    _add_text(
        slide, "Aegis ATV  사용 설명서",
        left=Inches(0.4), top=Inches(7.22), width=Inches(6), height=Inches(0.28),
        size=10, color=ICE,
    )
    _add_text(
        slide, f"{page_num} / {total}",
        left=Inches(12.0), top=Inches(7.22), width=Inches(1.0), height=Inches(0.28),
        size=10, color=ICE, align="right",
    )


def _card(slide, *, left, top, width, height, fill=SOFT_WHITE, border=None):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    card.adjustments[0] = 0.08
    _fill(card, fill)
    if border:
        _outline(card, border, 1.0)
    else:
        card.line.fill.background()
    return card


def _eyebrow(slide, label: str, *, top: Inches = DEFAULT_EYEBROW_TOP) -> None:
    """Small uppercase navy label above the title — PitchDeck pattern."""
    _add_text(
        slide, label.upper(),
        left=Inches(0.6), top=top, width=Inches(8.0), height=Inches(0.3),
        size=11, bold=True, color=CORAL,
    )


def _title(slide, text: str, *, top: Inches = DEFAULT_TITLE_TOP) -> None:
    _add_text(
        slide, text,
        left=Inches(0.6), top=top, width=Inches(12.1), height=Inches(1.0),
        size=32, bold=True, color=NAVY, font=HEAD_FONT,
    )


def _divider(slide, *, top: Inches = DEFAULT_DIVIDER_TOP) -> None:
    """Slim horizontal accent line under the title."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(0.6), Inches(0.04),
    )
    _fill(line, CORAL)


# ── slide builders ──────────────────────────────────────────────


def slide_01_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg(slide, NAVY)

    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(3.0),
        Inches(0.8), Inches(0.08),
    )
    _fill(bar, CORAL)

    _add_text(
        slide, "AEGIS ATV",
        left=Inches(0.6), top=Inches(3.2), width=Inches(12.0),
        height=Inches(0.7),
        size=18, bold=True, color=ICE, font=HEAD_FONT,
    )
    _add_text(
        slide, "사용 설명서",
        left=Inches(0.6), top=Inches(3.7), width=Inches(12.0),
        height=Inches(1.2),
        size=60, bold=True, color=WHITE, font=HEAD_FONT,
    )
    _add_text(
        slide,
        "AI 에이전트의 모든 행동을 실행 직전 검증하고 위조 불가능한 기록으로 남깁니다",
        left=Inches(0.6), top=Inches(4.95), width=Inches(12.0),
        height=Inches(0.5),
        size=20, color=ICE, font=HEAD_FONT,
    )
    _add_text(
        slide, "AegisData  ·  2026-05-15  ·  비전문가용 통합 가이드",
        left=Inches(0.6), top=Inches(6.7), width=Inches(12.0),
        height=Inches(0.4),
        size=12, color=ICE,
    )


def slide_02_what_is_it(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _eyebrow(slide, "한 페이지 요약")
    _title(slide, "Aegis 는 무엇을 하는 도구입니까?")
    _divider(slide)

    _add_text(
        slide,
        "AI 에이전트가 실수 또는 공격으로 시스템을 망가뜨리지 못하게 막고,\n"
        "동시에 모든 행동을 위조 불가능한 기록으로 남기는 도구입니다.",
        left=Inches(0.6), top=Inches(1.85), width=Inches(12.0), height=Inches(1.2),
        size=20, color=DARK_INK, font=BODY_FONT,
    )

    # 3 비유 카드
    items = [
        ("🚪", "자물쇠", "위험한 명령을 실행 직전에 막습니다"),
        ("📹", "CCTV", "모든 행동을 시간 순서대로 기록합니다"),
        ("🧾", "공증된 영수증", "암호로 서명되어 위조 불가능합니다"),
    ]
    card_w = Inches(3.95)
    card_h = Inches(2.4)
    gap = Inches(0.15)
    start_left = Inches(0.6)
    for i, (icon, title, desc) in enumerate(items):
        left = start_left + (card_w + gap) * i
        _card(slide, left=left, top=Inches(3.4), width=card_w, height=card_h)
        _add_text(
            slide, icon,
            left=left + Inches(0.2), top=Inches(3.5),
            width=Inches(1), height=Inches(0.7),
            size=36,
        )
        _add_text(
            slide, title,
            left=left + Inches(0.3), top=Inches(4.4),
            width=card_w - Inches(0.6), height=Inches(0.5),
            size=22, bold=True, color=NAVY,
        )
        _add_text(
            slide, desc,
            left=left + Inches(0.3), top=Inches(4.95),
            width=card_w - Inches(0.6), height=Inches(1.2),
            size=14, color=MUTED,
        )

    _footer(slide, "2")


def slide_03_chokepoint(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _eyebrow(slide, "어디에 끼어드나요?")
    _title(slide, '"below the model" — 결정과 실행 사이의 chokepoint')
    _divider(slide)

    # 3-layer column figure
    col_w = Inches(3.6)
    col_h = Inches(1.2)
    col_top = Inches(2.4)
    layers = [
        ("AI 에이전트", "Claude Code · OpenClaw · Codex · custom", ICE, NAVY),
        ("★ AEGIS", "16-step firewall · cryptographic audit", CORAL, WHITE),
        ("실제 도구", "셸 · DB · API · 결제 · 파일시스템", ICE, NAVY),
    ]
    for i, (label, sub, fill, fg) in enumerate(layers):
        left = Inches(0.6) + (col_w + Inches(0.5)) * i
        card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, col_top, col_w, col_h,
        )
        _fill(card, fill)
        _outline(card, NAVY, 1.0)
        _add_text(
            slide, label,
            left=left, top=col_top + Inches(0.18),
            width=col_w, height=Inches(0.5),
            size=24, bold=True, color=fg, align="center",
        )
        _add_text(
            slide, sub,
            left=left, top=col_top + Inches(0.7),
            width=col_w, height=Inches(0.4),
            size=12, color=fg if fill == CORAL else MUTED, align="center",
        )

    # arrows between layers
    for i in range(2):
        ax = Inches(0.6) + col_w + Inches(0.05) + (col_w + Inches(0.5)) * i
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, ax, col_top + Inches(0.45),
            Inches(0.4), Inches(0.3),
        )
        _fill(arrow, NAVY)

    # bottom annotation
    _add_text(
        slide,
        "모델의 안전 응답 필터보다 한 단계 더 아래에서 작동합니다.",
        left=Inches(0.6), top=Inches(4.4), width=Inches(12.0), height=Inches(0.4),
        size=16, color=DARK_INK, italic=True,
    )

    # 3 quick stats
    stats = [
        ("< 50ms", "정책 결정 지연 (p95)"),
        ("≥ 90%", "공격 차단율"),
        ("1 cmd", "외부 감사 가능"),
    ]
    sw = Inches(3.95)
    sh = Inches(1.5)
    s_top = Inches(5.1)
    for i, (big, sub) in enumerate(stats):
        left = Inches(0.6) + (sw + Inches(0.15)) * i
        _card(slide, left=left, top=s_top, width=sw, height=sh)
        _add_text(
            slide, big,
            left=left, top=s_top + Inches(0.15),
            width=sw, height=Inches(0.7),
            size=36, bold=True, color=NAVY, align="center", font=HEAD_FONT,
        )
        _add_text(
            slide, sub,
            left=left, top=s_top + Inches(0.95),
            width=sw, height=Inches(0.45),
            size=12, color=MUTED, align="center",
        )

    _footer(slide, "3")


def slide_04_personas(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _eyebrow(slide, "어떤 사람에게 필요한가")
    _title(slide, "5 가지 사용자 페르소나")
    _divider(slide)

    rows = [
        ("🧑‍💻", "Claude Code 일상 사용 개발자",
            "위험한 파일 삭제 자동 차단 + 모든 작업 자동 기록"),
        ("🏥", "병원 · 금융 · 정부 규제 산업",
            "EU AI Act · HIPAA · SOC 2 가 요구하는 변조 불가 감사 로그"),
        ("🤖", "AI 에이전트를 만드는 개발자",
            "OpenClaw · 자체 프레임워크에 보안 + 감사 layer 한 줄로"),
        ("🛡️", "기업 보안 · 컴플라이언스 팀",
            "여러 AI 도구 × 여러 LLM provider 의 행동을 하나의 대시보드"),
        ("🚀", "Multi-LLM 환경 운영자",
            "OpenRouter 환경에서 어느 provider 가 더 위험한지 정량 측정"),
    ]
    row_top = Inches(2.0)
    row_h = Inches(0.92)
    for i, (icon, who, what) in enumerate(rows):
        top = row_top + row_h * i
        # Icon circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.65), top + Inches(0.13),
            Inches(0.65), Inches(0.65),
        )
        _fill(circle, NAVY)
        _add_text(
            slide, icon,
            left=Inches(0.65), top=top + Inches(0.18),
            width=Inches(0.65), height=Inches(0.55),
            size=22, color=WHITE, align="center",
        )
        # Persona + benefit
        _add_text(
            slide, who,
            left=Inches(1.55), top=top + Inches(0.05),
            width=Inches(5.0), height=Inches(0.4),
            size=15, bold=True, color=NAVY,
        )
        _add_text(
            slide, what,
            left=Inches(1.55), top=top + Inches(0.42),
            width=Inches(11.0), height=Inches(0.5),
            size=13, color=DARK_INK,
        )

    _footer(slide, "4")


def slide_05_scenario(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _eyebrow(slide, "한 시나리오로 보는 동작")
    _title(slide, "삭제 사고 방지 — Aegis 없을 때 vs 있을 때")
    _divider(slide)

    # Left card: 없을 때
    _card(
        slide, left=Inches(0.6), top=Inches(1.9),
        width=Inches(6.0), height=Inches(4.6),
        fill=RGBColor(0xFD, 0xE9, 0xEA),
    )
    _add_text(
        slide, "❌  Aegis 없이",
        left=Inches(0.8), top=Inches(2.05),
        width=Inches(5.6), height=Inches(0.5),
        size=20, bold=True, color=CORAL,
    )
    _add_text(
        slide,
        '사용자: "tmp 정리해줘"\n'
        "Claude:  → 시스템 폴더 대상 재귀 삭제 명령\n"
        "        → 실행됨\n\n"
        "→ 시스템 망가짐 ❌\n"
        "→ 사후 분석 자료 없음 ❌",
        left=Inches(0.8), top=Inches(2.7),
        width=Inches(5.6), height=Inches(3.6),
        size=14, color=DARK_INK, font=MONO_FONT,
    )

    # Right card: 있을 때
    _card(
        slide, left=Inches(6.85), top=Inches(1.9),
        width=Inches(6.0), height=Inches(4.6),
        fill=RGBColor(0xE6, 0xF3, 0xEE),
    )
    _add_text(
        slide, "✓  Aegis 와 함께",
        left=Inches(7.05), top=Inches(2.05),
        width=Inches(5.6), height=Inches(0.5),
        size=20, bold=True, color=RGBColor(0x14, 0x77, 0x4A),
    )
    _add_text(
        slide,
        '사용자: "tmp 정리해줘"\n'
        "Claude:  → 같은 위험 명령 시도\n"
        "Aegis:   ⛔ BLOCK trace=abc123 (45ms)\n"
        "         reason: dangerous pattern\n"
        "         advise: security-reviewer\n\n"
        "→ 안 망가짐 ✓\n"
        "→ audit log 에 시도 영구 기록 ✓",
        left=Inches(7.05), top=Inches(2.7),
        width=Inches(5.6), height=Inches(3.6),
        size=14, color=DARK_INK, font=MONO_FONT,
    )

    _footer(slide, "5")


def slide_06_install(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _eyebrow(slide, "5분 설치")
    _title(slide, "4 가지 설치 옵션")
    _divider(slide)

    options = [
        ("A · Solo 개발자",
            "git clone https://github.com/happyikas/Aegis-ATV.git\n"
            "cd Aegis-ATV && uv sync\n"
            "uv run aegis install --mode local"),
        ("B · 한 줄 설치",
            "curl -LsSf https://raw.githubusercontent.com/\n"
            "  happyikas/Aegis-ATV/main/scripts/install.sh | bash"),
        ("C · Mac (Homebrew)",
            "brew tap happyikas/aegis \\\n"
            "  https://github.com/happyikas/Aegis-ATV.git\n"
            "brew install happyikas/aegis/aegis\n"
            "aegis install --mode local"),
        ("D · 기업 · 멀티 사용자",
            "docker compose up -d\n"
            "# localhost:8000 FastAPI 가동\n"
            "# 멀티 클라이언트가 이 service 검증"),
    ]
    card_w = Inches(6.05)
    card_h = Inches(2.45)
    start_left = Inches(0.6)
    start_top = Inches(1.9)
    for i, (title, cmd) in enumerate(options):
        row = i // 2
        col = i % 2
        left = start_left + (card_w + Inches(0.2)) * col
        top = start_top + (card_h + Inches(0.2)) * row
        _card(slide, left=left, top=top, width=card_w, height=card_h)
        _add_text(
            slide, title,
            left=left + Inches(0.25), top=top + Inches(0.18),
            width=card_w - Inches(0.5), height=Inches(0.45),
            size=16, bold=True, color=NAVY,
        )
        _add_text(
            slide, cmd,
            left=left + Inches(0.25), top=top + Inches(0.75),
            width=card_w - Inches(0.5), height=card_h - Inches(1.0),
            size=12, color=DARK_INK, font=MONO_FONT,
        )

    _footer(slide, "6")


def slide_07_commands(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _eyebrow(slide, "처음 5 가지 명령어")
    _title(slide, "설치 후 가장 먼저 익혀야 할 명령")
    _divider(slide)

    rows = [
        ("aegis status",         "설치 상태 + 운영 통계",      "매일 첫 명령"),
        ("aegis report",         "최근 24h 위험 요약 (5줄)",   "매일 / 매주"),
        ("aegis verify-audit",   "감사 체인 무결성 검증 (1s)", "사건 의심 / 매주"),
        ("aegis forensic last",  "최근 BLOCK 케이스 분석",     "왜 차단됐지?"),
        ("aegis advise",         "advisor 권고 종합",          "운영 개선"),
    ]

    # Header row
    hdr_top = Inches(1.9)
    _add_text(slide, "명령",       left=Inches(0.65),  top=hdr_top,
              width=Inches(3.5), height=Inches(0.35), size=12, bold=True, color=MUTED)
    _add_text(slide, "설명",       left=Inches(4.4),   top=hdr_top,
              width=Inches(5.0), height=Inches(0.35), size=12, bold=True, color=MUTED)
    _add_text(slide, "언제 쓰나",  left=Inches(9.6),   top=hdr_top,
              width=Inches(3.5), height=Inches(0.35), size=12, bold=True, color=MUTED)

    for i, (cmd, desc, when) in enumerate(rows):
        top = Inches(2.3) + Inches(0.55) * i
        # alternating background
        if i % 2 == 0:
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.6), top,
                Inches(12.13), Inches(0.5),
            )
            _fill(bg, SOFT_WHITE)
        _add_text(slide, cmd, left=Inches(0.7), top=top + Inches(0.08),
                  width=Inches(3.6), height=Inches(0.4), size=14,
                  font=MONO_FONT, color=NAVY, bold=True)
        _add_text(slide, desc, left=Inches(4.4), top=top + Inches(0.08),
                  width=Inches(5.0), height=Inches(0.4), size=13,
                  color=DARK_INK)
        _add_text(slide, when, left=Inches(9.6), top=top + Inches(0.08),
                  width=Inches(3.3), height=Inches(0.4), size=13,
                  color=MUTED, italic=True)

    # Sample output preview
    sample_top = Inches(5.4)
    sample = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), sample_top,
        Inches(12.13), Inches(1.65),
    )
    _fill(sample, DARK_INK)
    _add_text(
        slide,
        "$ aegis report\n"
        "  Calls: 1,243   ALLOW 1,198 · APPROVAL 38 · BLOCK 7\n"
        "  Top risks: 4× destructive  ·  2× credential-leak  ·  1× cost spike\n"
        "  Audit chain: ✓ intact, 12 sessions",
        left=Inches(0.85), top=sample_top + Inches(0.18),
        width=Inches(11.5), height=Inches(1.4),
        size=12, color=ICE, font=MONO_FONT,
    )

    _footer(slide, "7")


def slide_08_three_features(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _eyebrow(slide, "3 가지 핵심 기능")
    _title(slide, "Coach · Live · Doctor — 학습 → 모니터링 → 진단")
    _divider(slide)

    features = [
        ("🏋️", "ATV Coach",
            "정상 패턴 학습",
            "5단계 × 4 phase 자동 학습\n"
            "Observation → Shadow →\nAssisted → Production"),
        ("📊", "ATV Live",
            "실시간 모니터링",
            "Cost · Performance · Security\nby agent / channel / provider\n실시간 알림 daemon"),
        ("🔧", "ATV Doctor",
            "사건 분석 · 진단",
            "BLOCK 케이스 재현\n8-advisor 권고\n시간 되돌리기 (rollback)"),
    ]
    card_w = Inches(3.95)
    card_h = Inches(4.65)
    start_left = Inches(0.6)
    for i, (icon, name, tag, body) in enumerate(features):
        left = start_left + (card_w + Inches(0.15)) * i
        _card(slide, left=left, top=Inches(1.95),
              width=card_w, height=card_h)
        # Icon + name
        _add_text(slide, icon,
                  left=left + Inches(0.3), top=Inches(2.15),
                  width=Inches(1), height=Inches(0.8), size=42)
        _add_text(slide, name,
                  left=left + Inches(0.3), top=Inches(3.05),
                  width=card_w - Inches(0.6), height=Inches(0.55),
                  size=22, bold=True, color=NAVY)
        _add_text(slide, tag,
                  left=left + Inches(0.3), top=Inches(3.65),
                  width=card_w - Inches(0.6), height=Inches(0.45),
                  size=14, color=CORAL, italic=True)
        # Divider
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left + Inches(0.3), Inches(4.2),
            Inches(0.5), Inches(0.03),
        )
        _fill(bar, NAVY)
        _add_text(slide, body,
                  left=left + Inches(0.3), top=Inches(4.4),
                  width=card_w - Inches(0.6), height=Inches(2.0),
                  size=13, color=DARK_INK)

    _footer(slide, "8")


def slide_09_five_techs(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _eyebrow(slide, "PitchDeck 의 5 기술")
    _title(slide, "ATV · ATMU · sLLM · Crypto-Sign · Burn-in")
    _divider(slide)

    techs = [
        ("ATV",         "Agent Telemetry Vector",
            "모든 행동을 2,048-D 벡터로 인코딩"),
        ("ATMU",        "Trust Mgmt Unit",
            "16 단계 firewall — 실행 직전 정책 게이트"),
        ("sLLM",        "Judgment Engine",
            "양자화 3B 로컬 — ambiguous 케이스 second opinion"),
        ("Crypto-Sign", "Tamper-Evident",
            "SHA3 + Ed25519 체인 — 1 명령 외부 감사"),
        ("Burn-in",     "Release Gate",
            "1k+ 적대 시나리오 재현 — CI 게이트"),
    ]
    item_w = Inches(2.42)
    item_h = Inches(3.4)
    start_left = Inches(0.6)
    for i, (name, tag, desc) in enumerate(techs):
        left = start_left + (item_w + Inches(0.07)) * i
        _card(slide, left=left, top=Inches(2.0),
              width=item_w, height=item_h)
        # numbered circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left + Inches(0.85), Inches(2.2),
            Inches(0.7), Inches(0.7),
        )
        _fill(circle, NAVY)
        _add_text(slide, str(i + 1),
                  left=left + Inches(0.85), top=Inches(2.28),
                  width=Inches(0.7), height=Inches(0.6),
                  size=22, bold=True, color=WHITE, align="center",
                  font=HEAD_FONT)
        _add_text(slide, name,
                  left=left + Inches(0.15), top=Inches(3.05),
                  width=item_w - Inches(0.3), height=Inches(0.45),
                  size=18, bold=True, color=NAVY, align="center",
                  font=HEAD_FONT)
        _add_text(slide, tag,
                  left=left + Inches(0.15), top=Inches(3.5),
                  width=item_w - Inches(0.3), height=Inches(0.4),
                  size=11, color=CORAL, align="center", italic=True)
        _add_text(slide, desc,
                  left=left + Inches(0.2), top=Inches(4.0),
                  width=item_w - Inches(0.4), height=Inches(1.5),
                  size=11, color=DARK_INK, align="center")

    # Bottom summary
    _add_text(
        slide,
        "표준 벡터 인코딩 → 정책 적용 → 로컬 AI 검토 → 서명 체인 기록 → 적대 시나리오 재검증",
        left=Inches(0.6), top=Inches(5.85),
        width=Inches(12.0), height=Inches(0.4),
        size=14, color=NAVY, italic=True, align="center",
    )

    _footer(slide, "9")


def slide_10_plans(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _eyebrow(slide, "요금제")
    _title(slide, "Solo Free · Pro · Team · Enterprise")
    _divider(slide)

    plans = [
        ("Solo Free", "무료 (영구)", "Apache-2.0",
            ["✓ 16-step Firewall",
             "✓ 감사 체인",
             "✗ 8-advisor",
             "✗ sLLM judge",
             "✗ Sidecar"]),
        ("Pro", "$19/월", "상용",
            ["✓ 16-step Firewall",
             "✓ 감사 체인",
             "✓ 8-advisor",
             "✓ sLLM judge",
             "✗ Sidecar"]),
        ("Team", "$39/seat/월", "상용",
            ["✓ 16-step Firewall",
             "✓ 감사 체인",
             "✓ 8-advisor",
             "✓ sLLM judge",
             "✓ Sidecar"]),
        ("Enterprise", "별도", "상용",
            ["✓ Team 전부",
             "✓ Haiku judge",
             "✓ 우선 지원 + SLA",
             "✓ 보안 컨설팅",
             "✓ 커스텀 통합"]),
    ]
    card_w = Inches(3.05)
    card_h = Inches(4.6)
    start_left = Inches(0.6)
    for i, (name, price, lic, features) in enumerate(plans):
        left = start_left + (card_w + Inches(0.07)) * i
        is_pro = (name == "Pro")
        _card(
            slide, left=left, top=Inches(1.95),
            width=card_w, height=card_h,
            fill=ICE if is_pro else SOFT_WHITE,
            border=CORAL if is_pro else None,
        )
        _add_text(slide, name,
                  left=left + Inches(0.2), top=Inches(2.1),
                  width=card_w - Inches(0.4), height=Inches(0.5),
                  size=20, bold=True, color=NAVY, align="center")
        _add_text(slide, price,
                  left=left + Inches(0.2), top=Inches(2.65),
                  width=card_w - Inches(0.4), height=Inches(0.5),
                  size=22, bold=True, color=CORAL if is_pro else NAVY,
                  align="center", font=HEAD_FONT)
        _add_text(slide, lic,
                  left=left + Inches(0.2), top=Inches(3.2),
                  width=card_w - Inches(0.4), height=Inches(0.35),
                  size=11, color=MUTED, align="center", italic=True)
        # Divider
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left + Inches(0.6), Inches(3.65),
            Inches(0.4), Inches(0.03),
        )
        _fill(bar, NAVY)
        # Features
        for j, feat in enumerate(features):
            _add_text(slide, feat,
                      left=left + Inches(0.3), top=Inches(3.85) + Inches(0.4) * j,
                      width=card_w - Inches(0.5), height=Inches(0.35),
                      size=12,
                      color=DARK_INK if feat.startswith("✓") else MUTED)

    _footer(slide, "10")


def slide_11_integrations(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _eyebrow(slide, "통합 시나리오")
    _title(slide, "어떤 도구 · 환경에서 쓰나")
    _divider(slide)

    items = [
        ("Claude Code",
            "가장 흔한 사용 환경",
            "uv run aegis install --mode local --profile pro",
            "→ ~/.claude/settings.json 자동 패치 + 모든 tool call 자동 검증"),
        ("OpenClaw",
            "멀티 채널 agent (Telegram · Discord · Slack)",
            "npm install @happyikas/openclaw-plugin-aegis",
            "→ before_tool_call 후크 + sidecar 자동 검증"),
        ("OpenRouter",
            "300+ models · 60+ providers",
            "from aegis.integrations.openrouter import canonical_provider",
            "→ provider-drift advisor — 어느 provider 가 더 위험한지 정량화"),
        ("Hermes",
            "Self-improving agent",
            "docs/integrations/hermes.md",
            "→ 외부 감사 layer — self-improvement 의 cross-check"),
    ]
    card_w = Inches(6.05)
    card_h = Inches(2.45)
    start_left = Inches(0.6)
    start_top = Inches(1.95)
    for i, (name, desc, cmd, then) in enumerate(items):
        row = i // 2
        col = i % 2
        left = start_left + (card_w + Inches(0.2)) * col
        top = start_top + (card_h + Inches(0.2)) * row
        _card(slide, left=left, top=top, width=card_w, height=card_h)
        _add_text(slide, name,
                  left=left + Inches(0.25), top=top + Inches(0.15),
                  width=card_w - Inches(0.5), height=Inches(0.45),
                  size=18, bold=True, color=NAVY)
        _add_text(slide, desc,
                  left=left + Inches(0.25), top=top + Inches(0.62),
                  width=card_w - Inches(0.5), height=Inches(0.35),
                  size=12, color=MUTED, italic=True)
        _add_text(slide, cmd,
                  left=left + Inches(0.25), top=top + Inches(1.05),
                  width=card_w - Inches(0.5), height=Inches(0.65),
                  size=11, color=DARK_INK, font=MONO_FONT)
        _add_text(slide, then,
                  left=left + Inches(0.25), top=top + Inches(1.75),
                  width=card_w - Inches(0.5), height=Inches(0.6),
                  size=12, color=NAVY)

    _footer(slide, "11")


def slide_12_faq(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _eyebrow(slide, "자주 묻는 질문")
    _title(slide, "비전문가가 가장 많이 묻는 4 가지")
    _divider(slide)

    items = [
        ("Q. 내 데이터가 외부로 나가나요?",
            "--mode local 사용 시 0 byte 외부 송신 없음. "
            "감사 로그도 ~/.aegis/audit.jsonl 본인 머신만."),
        ("Q. 너무 느리지 않을까요?",
            "Solo Free 평균 5ms (dummy 룰). "
            "Pro/Team 평균 50ms (sLLM 포함, p95)."),
        ("Q. Claude 가 이미 안전 응답을 하는데 또 필요한가요?",
            "Claude refuse 는 모델 출력 수준. Aegis 는 출력 이후 "
            "실제 tool 실행 직전에 한 번 더 검증 — below the model."),
        ("Q. 감사 로그가 진짜 위조 불가능한가요?",
            "네. SHA3 hash chain + Ed25519 서명. "
            "aegis verify-audit 한 명령으로 1초 안에 외부 검증 가능."),
    ]
    row_h = Inches(1.15)
    start_top = Inches(1.95)
    for i, (q, a) in enumerate(items):
        top = start_top + row_h * i
        # accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), top + Inches(0.1),
            Inches(0.06), Inches(0.95),
        )
        _fill(bar, CORAL)
        _add_text(slide, q,
                  left=Inches(0.85), top=top + Inches(0.05),
                  width=Inches(11.8), height=Inches(0.5),
                  size=15, bold=True, color=NAVY)
        _add_text(slide, a,
                  left=Inches(0.85), top=top + Inches(0.5),
                  width=Inches(11.8), height=Inches(0.6),
                  size=13, color=DARK_INK)

    _footer(slide, "12")


def slide_13_troubleshooting(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _eyebrow(slide, "자주 발생하는 문제")
    _title(slide, "증상 → 원인 → 해결")
    _divider(slide)

    rows = [
        ("Claude Code 가 후크 무시",
            "Claude Code 재시작 필수. aegis status 로 settings.json 패치 상태 확인."),
        ("aegis verify-audit 실패",
            "감사 로그 수정 가능성. aegis forensic last → backup 으로 복원."),
        ("BLOCK 메시지 너무 많음",
            "Coach baseline 학습 부족. 1주일 AEGIS_BURNIN_SHADOW=1 후 train-m13."),
        ("Pro 활성 후에도 dummy judge",
            "echo $AEGIS_JUDGE_PROVIDER 확인. unset 또는 local-phi / haiku 설정."),
        ("OpenRouter route 가 (no-provider) 로",
            "canonical_provider() 헬퍼로 provider 문자열 생성하여 ATV header 에 stamp."),
        ("한국어 메시지 깨짐 (Windows)",
            "터미널 UTF-8 설정: chcp 65001."),
    ]
    # Header row
    hdr_top = Inches(1.95)
    _add_text(slide, "증상",  left=Inches(0.65), top=hdr_top,
              width=Inches(5.0), height=Inches(0.35), size=12, bold=True, color=MUTED)
    _add_text(slide, "해결",  left=Inches(5.9),  top=hdr_top,
              width=Inches(7.0), height=Inches(0.35), size=12, bold=True, color=MUTED)
    for i, (sym, sol) in enumerate(rows):
        top = Inches(2.35) + Inches(0.65) * i
        if i % 2 == 0:
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.6), top,
                Inches(12.13), Inches(0.6),
            )
            _fill(bg, SOFT_WHITE)
        _add_text(slide, sym, left=Inches(0.7), top=top + Inches(0.12),
                  width=Inches(5.1), height=Inches(0.45), size=13,
                  color=NAVY, bold=True)
        _add_text(slide, sol, left=Inches(5.9), top=top + Inches(0.12),
                  width=Inches(7.0), height=Inches(0.45), size=12,
                  color=DARK_INK)

    _footer(slide, "13")


def slide_14_context_memory(prs: Presentation) -> None:
    """ContextMemory + aegis doctor — PitchDeck "HARDWARE NEXT" 매핑."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _eyebrow(slide, "PitchDeck \"HARDWARE NEXT\" 매핑")
    _title(slide, "ContextMemory + aegis doctor")
    _divider(slide)

    # Subtitle
    _add_text(
        slide,
        "ATV 가 생성될 때마다 자동 저장 — CXL SSD / Computational SSD 의 software emulation",
        left=Inches(0.6), top=Inches(1.85),
        width=Inches(12.13), height=Inches(0.4),
        size=14, color=MUTED, italic=True,
    )

    # Left card: 두 store 의 역할 분리
    _card(
        slide, left=Inches(0.6), top=Inches(2.45),
        width=Inches(6.05), height=Inches(2.5),
    )
    _add_text(
        slide, "두 store 의 역할 분리",
        left=Inches(0.85), top=Inches(2.6),
        width=Inches(5.55), height=Inches(0.4),
        size=16, bold=True, color=NAVY,
    )
    _add_text(
        slide,
        "audit.jsonl\n"
        "  SHA3 + Ed25519 체인 (변조 증거)\n\n"
        "context_memory.jsonl  ← 신규\n"
        "  분석 fast-path. silicon 이행 시\n"
        "  near-storage compute 의 입력",
        left=Inches(0.85), top=Inches(3.05),
        width=Inches(5.55), height=Inches(1.9),
        size=13, color=DARK_INK, font=MONO_FONT,
    )

    # Right card: aegis doctor sample
    _card(
        slide, left=Inches(6.85), top=Inches(2.45),
        width=Inches(5.88), height=Inches(2.5),
    )
    _add_text(
        slide, "$ aegis doctor",
        left=Inches(7.1), top=Inches(2.6),
        width=Inches(5.4), height=Inches(0.4),
        size=14, bold=True, color=NAVY, font=MONO_FONT,
    )
    _add_text(
        slide,
        "📊 요약 — ALLOW 96.4% / BLOCK 0.5%\n"
        "💰 Cost  — $4.18 (94% Claude)\n"
        "⚡ Perf  — p95 47ms ✓\n"
        "🛡️ Sec  — step310 BLOCK 57%\n\n"
        "→ markdown 리포트 + 휴리스틱 권고",
        left=Inches(7.1), top=Inches(3.05),
        width=Inches(5.4), height=Inches(1.9),
        size=12, color=DARK_INK, font=MONO_FONT,
    )

    # Bottom strip: 3-tier value
    items = [
        ("매 tool call 자동 기록", "ATV 생성 시점에 두 파일 동시 append (defensive)"),
        ("3축 분석 + 권고", "Cost · Performance · Security 휴리스틱 advisor 11종 룰"),
        ("silicon-ready", "Same schema = CXL/CSD 의 spec. host-Python 이 후일 in-storage 로 이행"),
    ]
    item_w = Inches(4.05)
    item_h = Inches(1.5)
    start_left = Inches(0.6)
    for i, (head, body) in enumerate(items):
        left = start_left + (item_w + Inches(0.06)) * i
        # numbered circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left, Inches(5.15),
            Inches(0.45), Inches(0.45),
        )
        _fill(circle, CORAL)
        _add_text(slide, str(i + 1),
                  left=left, top=Inches(5.21),
                  width=Inches(0.45), height=Inches(0.4),
                  size=14, bold=True, color=WHITE, align="center")
        _add_text(slide, head,
                  left=left + Inches(0.55), top=Inches(5.18),
                  width=item_w - Inches(0.6), height=Inches(0.4),
                  size=13, bold=True, color=NAVY)
        _add_text(slide, body,
                  left=left, top=Inches(5.7),
                  width=item_w - Inches(0.1), height=item_h - Inches(0.6),
                  size=11, color=DARK_INK)

    _footer(slide, "14")


def slide_15_summary(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, NAVY)

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.8),
        Inches(0.8), Inches(0.08),
    )
    _fill(bar, CORAL)

    _add_text(slide, "한 문장 요약",
              left=Inches(0.6), top=Inches(3.0),
              width=Inches(12.0), height=Inches(0.5),
              size=14, bold=True, color=CORAL,
              font=HEAD_FONT)
    _add_text(
        slide,
        '"Aegis 는 당신의 AI 에이전트가 만드는 모든 행동을,\n'
        '실행 직전에 검증하고, 영원히 위조 불가능한 기록으로\n'
        '남기는 도구입니다."',
        left=Inches(0.6), top=Inches(3.4),
        width=Inches(12.0), height=Inches(2.2),
        size=28, bold=True, color=WHITE, font=HEAD_FONT,
    )
    _add_text(
        slide,
        "무료로 시작 · 노트북 밖으로 데이터 안 나감 · 한 명령으로 감사 가능",
        left=Inches(0.6), top=Inches(5.6),
        width=Inches(12.0), height=Inches(0.5),
        size=16, color=ICE, italic=True,
    )

    # install command box
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2.0), Inches(6.2),
        Inches(9.33), Inches(0.6),
    )
    _fill(box, DARK_INK)
    _outline(box, CORAL, 1.0)
    _add_text(
        slide, "$ uv run aegis install --mode local",
        left=Inches(2.0), top=Inches(6.3),
        width=Inches(9.33), height=Inches(0.4),
        size=18, color=ICE, font=MONO_FONT, align="center",
    )

    _add_text(
        slide,
        "GitHub: github.com/happyikas/Aegis-ATV  ·  aegisdata.ai  ·  datamonster@aegisdata.ai",
        left=Inches(0.6), top=Inches(7.15),
        width=Inches(12.13), height=Inches(0.3),
        size=10, color=ICE, align="center",
    )


# ── main ────────────────────────────────────────────────────────


def build(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_cover(prs)
    slide_02_what_is_it(prs)
    slide_03_chokepoint(prs)
    slide_04_personas(prs)
    slide_05_scenario(prs)
    slide_06_install(prs)
    slide_07_commands(prs)
    slide_08_three_features(prs)
    slide_09_five_techs(prs)
    slide_10_plans(prs)
    slide_11_integrations(prs)
    slide_12_faq(prs)
    slide_13_troubleshooting(prs)
    slide_14_context_memory(prs)
    slide_15_summary(prs)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"✓ wrote {out_path}  ({out_path.stat().st_size:,} bytes)")


if __name__ == "__main__":
    here = Path(__file__).resolve().parent.parent
    build(here / "docs" / "exports" / "USER_GUIDE.ko.pptx")
